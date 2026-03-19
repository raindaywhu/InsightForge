"""
PPT Generator Module
Generates .pptx files from JSON structure using python-pptx.
"""

import json
import tempfile
import urllib.request
import urllib.parse
from pathlib import Path
from typing import Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION


def _fix_unescaped_quotes(json_str: str) -> str:
    """
    Fix unescaped quotes inside JSON string values.
    This is a heuristic approach - we try to identify quotes that are likely
    inside string values (not JSON delimiters) and escape them.
    """
    result = []
    in_string = False
    escape_next = False
    
    i = 0
    while i < len(json_str):
        c = json_str[i]
        
        if escape_next:
            result.append(c)
            escape_next = False
            i += 1
            continue
            
        if c == '\\':
            result.append(c)
            escape_next = True
            i += 1
            continue
            
        if c == '"':
            if not in_string:
                # Starting a string
                in_string = True
                result.append(c)
            else:
                # We're in a string and see a quote
                # Check if this looks like an internal quote or string end
                # Look ahead: if next non-space char is : or , or } or ], it's likely string end
                rest = json_str[i+1:].lstrip()
                if rest and rest[0] in ':,}]':
                    # This is likely the end of the string
                    in_string = False
                    result.append(c)
                else:
                    # This is likely an internal quote - escape it
                    result.append('\\"')
            i += 1
            continue
        
        result.append(c)
        i += 1
    
    return ''.join(result)


# Theme configurations
THEMES = {
    "tech_blue": {
        "name": "科技蓝",
        "title_bg": RGBColor(26, 26, 46),      # Deep navy
        "accent": RGBColor(233, 69, 96),       # Coral red
        "accent2": RGBColor(52, 152, 219),     # Light blue
        "text_light": RGBColor(255, 255, 255),  # White
        "text_dark": RGBColor(51, 51, 51),      # Dark gray
        "content_bg": RGBColor(248, 249, 250),  # Light gray
        "header_bg": RGBColor(52, 73, 94),      # Steel blue
        "gradient_start": RGBColor(26, 26, 46),  # For gradient
        "gradient_end": RGBColor(44, 62, 80),
        "card_bg": RGBColor(255, 255, 255),    # Card background
        "border": RGBColor(220, 220, 220),     # Border color
    },
    "business_gray": {
        "name": "商务灰",
        "title_bg": RGBColor(45, 52, 54),       # Charcoal
        "accent": RGBColor(0, 184, 148),        # Teal
        "accent2": RGBColor(116, 185, 255),     # Light blue
        "text_light": RGBColor(255, 255, 255),
        "text_dark": RGBColor(45, 52, 54),
        "content_bg": RGBColor(245, 246, 250),
        "header_bg": RGBColor(99, 110, 114),
        "gradient_start": RGBColor(45, 52, 54),
        "gradient_end": RGBColor(72, 84, 96),
        "card_bg": RGBColor(255, 255, 255),
        "border": RGBColor(220, 220, 220),
    },
    "minimal_white": {
        "name": "简约白",
        "title_bg": RGBColor(255, 255, 255),    # White
        "accent": RGBColor(0, 123, 255),         # Blue
        "accent2": RGBColor(108, 117, 125),     # Gray
        "text_light": RGBColor(51, 51, 51),      # Dark for white bg
        "text_dark": RGBColor(51, 51, 51),
        "content_bg": RGBColor(255, 255, 255),
        "header_bg": RGBColor(240, 240, 240),
        "gradient_start": RGBColor(240, 240, 240),
        "gradient_end": RGBColor(255, 255, 255),
        "card_bg": RGBColor(250, 250, 250),
        "border": RGBColor(230, 230, 230),
    },
    "nature_green": {
        "name": "自然绿",
        "title_bg": RGBColor(39, 174, 96),       # Green
        "accent": RGBColor(243, 156, 18),        # Orange
        "accent2": RGBColor(46, 204, 113),       # Light green
        "text_light": RGBColor(255, 255, 255),
        "text_dark": RGBColor(45, 52, 54),
        "content_bg": RGBColor(245, 250, 245),
        "header_bg": RGBColor(46, 204, 113),
        "gradient_start": RGBColor(39, 174, 96),
        "gradient_end": RGBColor(46, 204, 113),
        "card_bg": RGBColor(255, 255, 255),
        "border": RGBColor(220, 230, 220),
    },
    "coral_sunset": {
        "name": "珊瑚夕阳",
        "title_bg": RGBColor(233, 69, 96),       # Coral
        "accent": RGBColor(255, 255, 255),       # White
        "accent2": RGBColor(255, 159, 67),       # Orange
        "text_light": RGBColor(255, 255, 255),
        "text_dark": RGBColor(51, 51, 51),
        "content_bg": RGBColor(255, 245, 245),
        "header_bg": RGBColor(214, 48, 71),
        "gradient_start": RGBColor(233, 69, 96),
        "gradient_end": RGBColor(255, 159, 67),
        "card_bg": RGBColor(255, 255, 255),
        "border": RGBColor(255, 220, 220),
    },
    "ocean_depth": {
        "name": "深海蓝",
        "title_bg": RGBColor(12, 35, 64),        # Deep ocean
        "accent": RGBColor(0, 210, 211),         # Cyan
        "accent2": RGBColor(72, 202, 228),       # Light cyan
        "text_light": RGBColor(255, 255, 255),
        "text_dark": RGBColor(51, 51, 51),
        "content_bg": RGBColor(240, 248, 255),
        "header_bg": RGBColor(20, 60, 90),
        "gradient_start": RGBColor(12, 35, 64),
        "gradient_end": RGBColor(20, 60, 100),
        "card_bg": RGBColor(255, 255, 255),
        "border": RGBColor(200, 220, 240),
    },
}


class PPTGenerator:
    """Generate PowerPoint presentations from JSON structure."""

    SLIDE_TYPES = {
        "title": "_create_title_slide",
        "agenda": "_create_content_slide",
        "key_findings": "_create_content_slide",
        "methodology": "_create_content_slide",
        "analysis": "_create_content_slide",
        "diagram": "_create_image_slide",  # diagram 使用图片幻灯片
        "chart": "_create_chart_slide",
        "data": "_create_chart_slide",  # data 类型也用图表
        "swot": "_create_swot_slide",  # SWOT 矩阵图
        "tech_curve": "_create_tech_curve_slide",  # 技术曲线
        "image": "_create_image_slide",  # 纯图片幻灯片
        "conclusion": "_create_content_slide",
        "recommendations": "_create_content_slide",
        "closing": "_create_closing_slide",
        "timeline": "_create_timeline_slide",  # 新增时间线
    }

    def __init__(self, theme: str = "tech_blue"):
        """Initialize PPT Generator with theme.
        
        Args:
            theme: Theme name (tech_blue, business_gray, minimal_white, nature_green)
        """
        self.theme_name = theme
        self.theme = THEMES.get(theme, THEMES["tech_blue"])
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9 Widescreen
        self.prs.slide_height = Inches(7.5)

    def generate(self, ppt_structure: dict | str, output_path: str | Path) -> str:
        """
        Generate PPT from JSON structure.

        Args:
            ppt_structure: JSON structure dict or path to JSON file
            output_path: Output path for the .pptx file

        Returns:
            Path to generated PPT file
        """
        # Load structure if path provided
        if isinstance(ppt_structure, str):
            with open(ppt_structure, "r", encoding="utf-8") as f:
                content = f.read()
            # Remove markdown code block markers if present
            if content.startswith("```"):
                first_newline = content.find("\n")
                if first_newline != -1:
                    content = content[first_newline + 1:]
            if content.rstrip().endswith("```"):
                content = content.rstrip()[:-3].rstrip()
            # Fix unescaped quotes in JSON
            content = _fix_unescaped_quotes(content)
            try:
                ppt_structure = json.loads(content)
            except json.JSONDecodeError as e:
                import re
                match = re.search(r'\{[\s\S]*\}', content)
                if match:
                    try:
                        ppt_structure = json.loads(match.group())
                    except json.JSONDecodeError as inner_e:
                        raise ValueError(f"Failed to parse JSON structure: {e}, inner: {inner_e}")
                else:
                    raise ValueError(f"No valid JSON found in file: {e}")

        # Validate structure
        if "slides" not in ppt_structure:
            raise ValueError("Invalid PPT structure: missing 'slides' key")

        # Create slides
        for slide_data in ppt_structure["slides"]:
            self._create_slide(slide_data)

        # Save presentation
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))

        return str(output_path)

    def _create_slide(self, slide_data: dict) -> None:
        """Create a single slide based on type."""
        slide_type = slide_data.get("type", "content")
        method_name = self.SLIDE_TYPES.get(slide_type, "_create_content_slide")
        method = getattr(self, method_name)
        method(slide_data)

    def _create_title_slide(self, slide_data: dict) -> None:
        """Create title slide (cover page) with enhanced visual design."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Add background shape
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.theme["title_bg"]
        shape.line.fill.background()

        # Add decorative accent bar at bottom
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(6.8), self.prs.slide_width, Inches(0.7)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.theme["accent"]
        accent_bar.line.fill.background()

        # Add decorative circle (top right)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(10.5), Inches(-0.5), Inches(3), Inches(3)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = self.theme["accent2"]
        # Set transparency by adjusting alpha
        try:
            fill = circle.fill
            fill.solid()
            # Add slight transparency effect by using a lighter shade
        except:
            pass
        circle.line.fill.background()

        # Add decorative line (left side)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(2.8), Inches(0.08), Inches(1.5)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.theme["accent"]
        line.line.fill.background()

        # Add title
        content = slide_data.get("content", [])
        if content:
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(2.5), Inches(11.0), Inches(1.5)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = content[0]
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = self.theme["text_light"]
            p.alignment = PP_ALIGN.LEFT

            # Subtitle
            if len(content) > 1:
                p = tf.add_paragraph()
                p.text = content[1]
                p.font.size = Pt(24)
                p.font.color.rgb = self.theme["accent2"]
                p.alignment = PP_ALIGN.LEFT
                p.space_before = Pt(16)

            # Footer / Date
            if len(content) > 2:
                footer_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(6.2), Inches(11.0), Inches(0.5)
                )
                tf = footer_box.text_frame
                p = tf.paragraphs[0]
                p.text = content[2]
                p.font.size = Pt(16)
                p.font.color.rgb = self.theme["text_light"]
                p.alignment = PP_ALIGN.LEFT

    def _create_content_slide(self, slide_data: dict) -> None:
        """Create content slide with enhanced visual design.
        
        Supports layouts:
        - default: header + bullet points
        - cards: content displayed as cards
        - two_column: two-column layout
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Check for layout type
        layout = slide_data.get("layout", "default")
        
        # Add header bar with accent line
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.theme["header_bg"]
        header.line.fill.background()

        # Add accent line below header
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(1.2), self.prs.slide_width, Inches(0.05)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = self.theme["accent"]
        accent_line.line.fill.background()

        # Add title with icon placeholder
        title = slide_data.get("title", "")
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.7)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = self.theme["text_light"]

        # Route to specific layout
        if layout == "cards":
            self._create_cards_layout(slide, slide_data)
        elif layout == "two_column":
            self._create_two_column_layout(slide, slide_data)
        else:
            self._create_default_layout(slide, slide_data)

        # Add slide number
        slide_num = slide_data.get("slide_number", "")
        if slide_num:
            num_box = slide.shapes.add_textbox(
                Inches(12.5), Inches(7), Inches(0.5), Inches(0.3)
            )
            tf = num_box.text_frame
            p = tf.paragraphs[0]
            p.text = str(slide_num)
            p.font.size = Pt(12)
            p.font.color.rgb = self.theme["text_dark"]
            p.alignment = PP_ALIGN.RIGHT

    def _create_default_layout(self, slide, slide_data: dict) -> None:
        """Create default bullet point layout with enhanced styling."""
        content = slide_data.get("content", [])
        if content:
            content_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(1.6), Inches(11.933), Inches(5.5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True

            for i, item in enumerate(content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                # Enhanced bullet styling
                p.text = "●  " + item
                p.font.size = Pt(20)
                p.font.color.rgb = self.theme["text_dark"]
                p.space_after = Pt(16)
                p.level = 0

    def _create_cards_layout(self, slide, slide_data: dict) -> None:
        """Create card-style layout for content.
        
        Expected format:
        {
            "layout": "cards",
            "cards": [
                {"title": "Card 1", "content": "Description"},
                {"title": "Card 2", "content": "Description"},
            ]
        }
        """
        cards = slide_data.get("cards", slide_data.get("content", []))
        
        if not cards:
            return
            
        # Determine card layout (2 or 3 columns)
        num_cards = len(cards)
        if num_cards <= 2:
            cols = 2
            card_width = 5.5
            card_height = 2.5
        else:
            cols = 3
            card_width = 3.8
            card_height = 2.5
            
        spacing = 0.5
        start_x = (13.333 - (cols * card_width + (cols - 1) * spacing)) / 2
        
        for i, card in enumerate(cards):
            if isinstance(card, str):
                card_title = ""
                card_content = card
            else:
                card_title = card.get("title", "")
                card_content = card.get("content", "")
            
            row = i // cols
            col = i % cols
            
            x = start_x + col * (card_width + spacing)
            y = 1.8 + row * (card_height + 0.3)
            
            # Create card background
            card_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), 
                Inches(card_width), Inches(card_height)
            )
            card_shape.fill.solid()
            card_shape.fill.fore_color.rgb = self.theme["card_bg"]
            card_shape.line.color.rgb = self.theme["border"]
            card_shape.line.width = Pt(1)
            
            # Add accent bar at top of card
            accent = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                Inches(card_width), Inches(0.08)
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = self.theme["accent"]
            accent.line.fill.background()
            
            # Add card title
            if card_title:
                title_box = slide.shapes.add_textbox(
                    Inches(x + 0.2), Inches(y + 0.25),
                    Inches(card_width - 0.4), Inches(0.5)
                )
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = card_title
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = self.theme["accent"]
            
            # Add card content
            content_y = y + 0.25 + (0.5 if card_title else 0)
            content_height = card_height - 0.5 - (0.5 if card_title else 0)
            
            content_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(content_y),
                Inches(card_width - 0.4), Inches(content_height)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = card_content
            p.font.size = Pt(14)
            p.font.color.rgb = self.theme["text_dark"]

    def _create_two_column_layout(self, slide, slide_data: dict) -> None:
        """Create two-column layout.
        
        Expected format:
        {
            "layout": "two_column",
            "left": {"title": "...", "content": [...]},
            "right": {"title": "...", "content": [...]}
        }
        """
        left = slide_data.get("left", {})
        right = slide_data.get("right", {})
        
        # Left column
        if left:
            self._create_column(slide, left, 0.5, 1.6, 5.8)
            
        # Right column
        if right:
            self._create_column(slide, right, 6.8, 1.6, 5.8)
            
        # Add vertical divider
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(6.4), Inches(1.8),
            Inches(0.03), Inches(5.0)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = self.theme["border"]
        divider.line.fill.background()

    def _create_column(self, slide, col_data: dict, x: float, y: float, width: float) -> None:
        """Create a single column with optional title and content."""
        title = col_data.get("title", "")
        content = col_data.get("content", [])
        
        current_y = y
        
        # Add column title
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(x), Inches(current_y), Inches(width), Inches(0.5)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.theme["accent"]
            current_y += 0.6
        
        # Add content
        if content:
            content_box = slide.shapes.add_textbox(
                Inches(x), Inches(current_y), Inches(width), Inches(4.5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for i, item in enumerate(content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = "●  " + item
                p.font.size = Pt(16)
                p.font.color.rgb = self.theme["text_dark"]
                p.space_after = Pt(10)

    def _create_timeline_slide(self, slide_data: dict) -> None:
        """Create timeline slide.
        
        Expected format:
        {
            "type": "timeline",
            "title": "Timeline",
            "items": [
                {"date": "2024 Q1", "title": "Event 1", "description": "..."},
                {"date": "2024 Q2", "title": "Event 2", "description": "..."},
            ]
        }
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Add header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.theme["header_bg"]
        header.line.fill.background()

        # Add accent line
        accent_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(1.2), self.prs.slide_width, Inches(0.05)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = self.theme["accent"]
        accent_line.line.fill.background()

        # Add title
        title = slide_data.get("title", "")
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.7)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = self.theme["text_light"]

        # Add timeline
        items = slide_data.get("items", [])
        if items:
            # Draw timeline line
            line_y = 3.5
            timeline_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(1), Inches(line_y),
                Inches(11.333), Inches(0.05)
            )
            timeline_line.fill.solid()
            timeline_line.fill.fore_color.rgb = self.theme["border"]
            timeline_line.line.fill.background()

            # Add items
            num_items = len(items)
            spacing = 11.333 / max(num_items, 1)
            
            for i, item in enumerate(items):
                x = 1 + i * spacing + spacing / 2
                
                # Add circle marker
                circle = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL, Inches(x - 0.15), Inches(line_y - 0.15),
                    Inches(0.35), Inches(0.35)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = self.theme["accent"]
                circle.line.fill.background()
                
                # Add date above line
                date = item.get("date", "")
                if date:
                    date_box = slide.shapes.add_textbox(
                        Inches(x - 0.8), Inches(line_y - 0.7),
                        Inches(1.6), Inches(0.4)
                    )
                    tf = date_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = date
                    p.font.size = Pt(12)
                    p.font.bold = True
                    p.font.color.rgb = self.theme["accent"]
                    p.alignment = PP_ALIGN.CENTER
                
                # Add title and description below line
                item_title = item.get("title", "")
                item_desc = item.get("description", "")
                
                text_box = slide.shapes.add_textbox(
                    Inches(x - 0.8), Inches(line_y + 0.4),
                    Inches(1.6), Inches(1.5)
                )
                tf = text_box.text_frame
                tf.word_wrap = True
                
                p = tf.paragraphs[0]
                p.text = item_title
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = self.theme["text_dark"]
                p.alignment = PP_ALIGN.CENTER
                
                if item_desc:
                    p = tf.add_paragraph()
                    p.text = item_desc
                    p.font.size = Pt(9)
                    p.font.color.rgb = self.theme["text_dark"]
                    p.alignment = PP_ALIGN.CENTER

    def _create_closing_slide(self, slide_data: dict) -> None:
        """Create closing slide (thank you page) with enhanced design."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Add background
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.theme["title_bg"]
        shape.line.fill.background()

        # Add decorative elements
        # Top right circle
        circle1 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(10.5), Inches(-0.5), Inches(3), Inches(3)
        )
        circle1.fill.solid()
        circle1.fill.fore_color.rgb = self.theme["accent2"]
        circle1.line.fill.background()

        # Bottom left circle
        circle2 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(-1), Inches(5), Inches(3), Inches(3)
        )
        circle2.fill.solid()
        circle2.fill.fore_color.rgb = self.theme["accent"]
        circle2.line.fill.background()

        # Add thank you text with accent line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(2.3), Inches(2.333), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.theme["accent"]
        line.line.fill.background()

        thank_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(2)
        )
        tf = thank_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Thank You"
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = self.theme["text_light"]
        p.alignment = PP_ALIGN.CENTER

        # Add contact info
        content = slide_data.get("content", [])
        if content:
            contact_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.5), Inches(12.333), Inches(1.5)
            )
            tf = contact_box.text_frame
            for i, item in enumerate(content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(18)
                p.font.color.rgb = self.theme["text_light"]
                p.alignment = PP_ALIGN.CENTER
                p.space_after = Pt(8)

    def _create_chart_slide(self, slide_data: dict) -> None:
        """Create slide with chart visualization.
        
        Expected slide_data format:
        {
            "type": "chart",
            "title": "Chart Title",
            "chart_type": "bar" | "column" | "line" | "pie",
            "chart_data": {
                "categories": ["Cat1", "Cat2", "Cat3"],
                "series": [
                    {"name": "Series 1", "values": [10, 20, 30]},
                    {"name": "Series 2", "values": [15, 25, 35]}
                ]
            },
            "content": ["Additional bullet points"]  # optional
        }
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Add header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.theme["header_bg"]
        header.line.fill.background()

        # Add title
        title = slide_data.get("title", "")
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.7)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = self.theme["text_light"]

        # Check if chart_data exists
        chart_data = slide_data.get("chart_data")
        if chart_data:
            chart_type = slide_data.get("chart_type", "column")
            self._add_chart(slide, chart_type, chart_data)

        # Add bullet points if provided (below chart or as main content)
        content = slide_data.get("content", [])
        if content:
            # Position below chart or full height if no chart
            top = Inches(4.8) if chart_data else Inches(1.5)
            height = Inches(2.2) if chart_data else Inches(5.5)
            
            content_box = slide.shapes.add_textbox(
                Inches(0.7), top, Inches(11.933), height
            )
            tf = content_box.text_frame
            tf.word_wrap = True

            for i, item in enumerate(content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = "\u2022 " + item
                p.font.size = Pt(16) if chart_data else Pt(20)
                p.font.color.rgb = self.theme["text_dark"]
                p.space_after = Pt(8)
                p.level = 0

    def _add_chart(self, slide, chart_type: str, chart_data: dict) -> None:
        """Add a chart to the slide.
        
        Args:
            slide: The slide to add chart to
            chart_type: Type of chart (bar, column, line, pie)
            chart_data: Chart data with categories and series
        """
        # Prepare chart data
        data = CategoryChartData()
        data.categories = chart_data.get("categories", [])
        
        series_list = chart_data.get("series", [])
        for series in series_list:
            data.add_series(series.get("name", ""), series.get("values", []))

        # Map chart type to python-pptx chart type
        chart_type_map = {
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
        }
        
        xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        
        # Add chart to slide
        x, y, cx, cy = Inches(0.7), Inches(1.5), Inches(11.933), Inches(3.2)
        chart = slide.shapes.add_chart(
            xl_chart_type, x, y, cx, cy, data
        ).chart

        # Style the chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        
        # Style plot area
        plot = chart.plots[0]
        plot.has_data_labels = True
        
        # For pie charts, show percentages
        if chart_type == "pie":
            plot.data_labels.show_percentage = True
            plot.data_labels.show_value = False
        else:
            plot.data_labels.show_value = True
            plot.data_labels.show_percentage = False

    def _create_swot_slide(self, slide_data: dict) -> None:
        """Create SWOT matrix visualization.
        
        Expected slide_data format:
        {
            "type": "swot",
            "title": "SWOT Analysis",
            "strengths": ["Item 1", "Item 2", ...],
            "weaknesses": ["Item 1", "Item 2", ...],
            "opportunities": ["Item 1", "Item 2", ...],
            "threats": ["Item 1", "Item 2", ...]
        }
        
        Or with content wrapper:
        {
            "type": "swot",
            "title": "SWOT Analysis",
            "content": {
                "strengths": [...],
                "weaknesses": [...],
                "opportunities": [...],
                "threats": [...]
            }
        }
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Add header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.theme["header_bg"]
        header.line.fill.background()

        # Add title
        title = slide_data.get("title", "SWOT Analysis")
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.7)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = self.theme["text_light"]

        # Get SWOT data
        content = slide_data.get("content", slide_data)
        strengths = content.get("strengths", [])
        weaknesses = content.get("weaknesses", [])
        opportunities = content.get("opportunities", [])
        threats = content.get("threats", [])

        # SWOT matrix colors and labels
        swot_config = [
            {
                "key": "strengths",
                "label": "Strengths\n优势",
                "items": strengths,
                "bg_color": RGBColor(46, 204, 113),  # Green
                "x": 0.5, "y": 1.5
            },
            {
                "key": "weaknesses", 
                "label": "Weaknesses\n劣势",
                "items": weaknesses,
                "bg_color": RGBColor(231, 76, 60),  # Red
                "x": 6.9, "y": 1.5
            },
            {
                "key": "opportunities",
                "label": "Opportunities\n机会",
                "items": opportunities,
                "bg_color": RGBColor(52, 152, 219),  # Blue
                "x": 0.5, "y": 4.5
            },
            {
                "key": "threats",
                "label": "Threats\n威胁",
                "items": threats,
                "bg_color": RGBColor(241, 196, 15),  # Yellow
                "x": 6.9, "y": 4.5
            }
        ]

        # Create SWOT quadrants
        box_width = 6.0
        box_height = 2.8

        for quad in swot_config:
            # Create quadrant background
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                Inches(quad["x"]), Inches(quad["y"]),
                Inches(box_width), Inches(box_height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = quad["bg_color"]
            shape.line.color.rgb = RGBColor(100, 100, 100)
            shape.line.width = Pt(1)

            # Add label
            label_box = slide.shapes.add_textbox(
                Inches(quad["x"] + 0.1), Inches(quad["y"] + 0.1),
                Inches(box_width - 0.2), Inches(0.6)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = quad["label"]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER

            # Add items
            if quad["items"]:
                items_box = slide.shapes.add_textbox(
                    Inches(quad["x"] + 0.2), Inches(quad["y"] + 0.75),
                    Inches(box_width - 0.4), Inches(box_height - 0.85)
                )
                tf = items_box.text_frame
                tf.word_wrap = True

                for i, item in enumerate(quad["items"][:5]):  # Max 5 items
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = Pt(12)
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.space_after = Pt(4)

        # Add axis labels
        # Internal/External label
        ie_label = slide.shapes.add_textbox(
            Inches(0.1), Inches(4.2), Inches(0.4), Inches(0.3)
        )
        tf = ie_label.text_frame
        p = tf.paragraphs[0]
        p.text = "内部"
        p.font.size = Pt(10)
        p.font.color.rgb = self.theme["text_dark"]
        
        ie_label2 = slide.shapes.add_textbox(
            Inches(0.1), Inches(5.8), Inches(0.4), Inches(0.3)
        )
        tf = ie_label2.text_frame
        p = tf.paragraphs[0]
        p.text = "外部"
        p.font.size = Pt(10)
        p.font.color.rgb = self.theme["text_dark"]

    def _create_tech_curve_slide(self, slide_data: dict) -> None:
        """Create technology curve visualization (similar to Gartner Hype Cycle).
        
        Expected slide_data format:
        {
            "type": "tech_curve",
            "title": "Technology Hype Cycle",
            "technologies": [
                {"name": "Tech 1", "position": 0.2, "peak": "2024"},
                {"name": "Tech 2", "position": 0.5, "peak": "2025"},
                {"name": "Tech 3", "position": 0.8, "peak": "2026"},
            ],
            "phases": [
                {"name": "Innovation Trigger", "range": [0, 0.2]},
                {"name": "Peak of Inflated Expectations", "range": [0.2, 0.4]},
                {"name": "Trough of Disillusionment", "range": [0.4, 0.6]},
                {"name": "Slope of Enlightenment", "range": [0.6, 0.8]},
                {"name": "Plateau of Productivity", "range": [0.8, 1.0]}
            ]
        }
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Add header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.theme["header_bg"]
        header.line.fill.background()

        # Add title
        title = slide_data.get("title", "Technology Hype Cycle")
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.7)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = self.theme["text_light"]

        # Curve parameters
        curve_left = 0.8
        curve_top = 1.8
        curve_width = 11.5
        curve_height = 4.0

        # Default phases if not provided
        phases = slide_data.get("phases", [
            {"name": "创新萌芽期", "range": [0, 0.2]},
            {"name": "期望膨胀期", "range": [0.2, 0.4]},
            {"name": "泡沫破裂期", "range": [0.4, 0.6]},
            {"name": "稳步爬升期", "range": [0.6, 0.8]},
            {"name": "生产成熟期", "range": [0.8, 1.0]}
        ])

        # Draw phase backgrounds
        phase_colors = [
            RGBColor(155, 89, 182),  # Purple
            RGBColor(231, 76, 60),   # Red
            RGBColor(241, 196, 15),  # Yellow
            RGBColor(46, 204, 113),  # Green
            RGBColor(52, 152, 219),  # Blue
        ]

        phase_width = curve_width / len(phases)
        for i, phase in enumerate(phases):
            # Phase background
            x = curve_left + i * phase_width
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(curve_top),
                Inches(phase_width), Inches(curve_height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = phase_colors[i % len(phase_colors)]
            # Set transparency (lighter colors)
            shape.fill.fore_color.brightness = 0.6
            shape.line.color.rgb = RGBColor(200, 200, 200)
            shape.line.width = Pt(0.5)

            # Phase label
            label_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(curve_top + curve_height + 0.1),
                Inches(phase_width - 0.2), Inches(0.6)
            )
            tf = label_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = phase["name"]
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.theme["text_dark"]
            p.alignment = PP_ALIGN.CENTER

        # Draw curve line (simplified as wavy line using shapes)
        # We'll use a series of connected lines
        curve_points = [
            (0.0, 0.5),   # Start low
            (0.15, 0.3),  # Rising
            (0.25, 0.1),  # Peak
            (0.35, 0.15), # Slight drop
            (0.45, 0.6),  # Trough
            (0.55, 0.7),  # Still low
            (0.65, 0.5),  # Rising again
            (0.75, 0.35), # Continues rising
            (0.85, 0.3),  # Levels off
            (1.0, 0.25),  # End
        ]

        # Draw curve as line (simplified - just draw the path)
        # Since python-pptx doesn't have freeform shapes easily, 
        # we'll draw a simplified curve using a shape
        curve_line = slide.shapes.add_shape(
            MSO_SHAPE.WAVE,
            Inches(curve_left), Inches(curve_top + 0.5),
            Inches(curve_width), Inches(curve_height - 1.0)
        )
        curve_line.fill.background()
        curve_line.line.color.rgb = self.theme["accent"]
        curve_line.line.width = Pt(3)

        # Add technology markers
        technologies = slide_data.get("technologies", [])
        for tech in technologies[:8]:  # Max 8 technologies
            pos = tech.get("position", 0.5)
            name = tech.get("name", "Technology")
            
            # Calculate position on curve
            x = curve_left + pos * curve_width
            
            # Find y position on curve (simplified)
            # Map position to curve height
            if pos < 0.3:
                y_offset = 0.3 + (0.3 - pos) * 1.5  # Peak area
            elif pos < 0.6:
                y_offset = 0.5 + (pos - 0.3) * 0.5  # Trough
            else:
                y_offset = 0.6 - (pos - 0.6) * 0.5  # Recovery
            
            y = curve_top + y_offset * curve_height

            # Draw marker
            marker = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x - 0.15), Inches(y - 0.1),
                Inches(0.3), Inches(0.2)
            )
            marker.fill.solid()
            marker.fill.fore_color.rgb = self.theme["accent"]
            marker.line.fill.background()

            # Add label
            label = slide.shapes.add_textbox(
                Inches(x - 0.5), Inches(y + 0.15),
                Inches(1.0), Inches(0.4)
            )
            tf = label.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = name
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = self.theme["text_dark"]
            p.alignment = PP_ALIGN.CENTER

        # Add Y-axis label
        y_label = slide.shapes.add_textbox(
            Inches(0.1), Inches(curve_top + 1.5),
            Inches(0.5), Inches(1.0)
        )
        tf = y_label.text_frame
        p = tf.paragraphs[0]
        p.text = "期望值"
        p.font.size = Pt(10)
        p.font.color.rgb = self.theme["text_dark"]

    def _create_image_slide(self, slide_data: dict) -> None:
        """Create slide with image.
        
        Expected slide_data format:
        {
            "type": "image",
            "title": "Image Title",
            "image": {
                "path": "/path/to/image.png",  # local path or URL
                "alt_text": "Image description",
                "position": "center" | "left" | "right",
                "width": 6.0,  # optional, in inches
                "height": 4.0  # optional, in inches
            },
            "content": ["Additional bullet points"]  # optional
        }
        
        Or simplified format:
        {
            "type": "image",
            "title": "Image Title",
            "image": "/path/to/image.png"  # just a path
        }
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Add header bar
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.theme["header_bg"]
        header.line.fill.background()

        # Add title
        title = slide_data.get("title", "")
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.7)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = self.theme["text_light"]

        # Process image
        image_data = slide_data.get("image")
        if image_data:
            image_path = None
            
            # Handle both string path and dict format
            if isinstance(image_data, str):
                image_path = image_data
                alt_text = ""
                position = "center"
                width = 8.0
                height = 4.0
            else:
                image_path = image_data.get("path", "")
                alt_text = image_data.get("alt_text", "")
                position = image_data.get("position", "center")
                width = image_data.get("width", 8.0)
                height = image_data.get("height", 4.0)
            
            # Download if URL
            local_path = self._resolve_image_path(image_path)
            
            if local_path and Path(local_path).exists():
                # Calculate position
                x = (13.333 - width) / 2  # center horizontally
                if position == "left":
                    x = 0.5
                elif position == "right":
                    x = 13.333 - width - 0.5
                
                y = Inches(1.5)
                
                try:
                    pic = slide.shapes.add_picture(
                        str(local_path), Inches(x), y, Inches(width), Inches(height)
                    )
                    if alt_text:
                        pic.name = alt_text
                except Exception as e:
                    # If image fails, add placeholder text
                    self._add_image_placeholder(slide, image_path, e)

        # Add bullet points if provided
        content = slide_data.get("content", [])
        if content:
            content_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(5.8), Inches(11.933), Inches(1.2)
            )
            tf = content_box.text_frame
            tf.word_wrap = True

            for i, item in enumerate(content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = "\u2022 " + item
                p.font.size = Pt(14)
                p.font.color.rgb = self.theme["text_dark"]
                p.space_after = Pt(6)

    def _resolve_image_path(self, path: str) -> str | None:
        """Resolve image path - download if URL, return local path.
        
        Args:
            path: Local path or URL
            
        Returns:
            Local file path or None if failed
        """
        if not path:
            return None
            
        # Check if URL
        if path.startswith(("http://", "https://")):
            try:
                # Download to temp file
                temp_dir = Path(tempfile.gettempdir()) / "insightforge_images"
                temp_dir.mkdir(parents=True, exist_ok=True)
                
                # Create filename from URL
                parsed = urllib.parse.urlparse(path)
                filename = Path(parsed.path).stem or "image"
                ext = Path(parsed.path).suffix or ".png"
                local_path = temp_dir / f"{filename}{ext}"
                
                # Download
                urllib.request.urlretrieve(path, str(local_path))
                return str(local_path)
            except Exception as e:
                print(f"[WARN] Failed to download image {path}: {e}")
                return None
        else:
            # Local path - check if exists
            if Path(path).exists():
                return path
            return None

    def _add_image_placeholder(self, slide, image_path: str, error: Exception) -> None:
        """Add placeholder text when image fails to load."""
        placeholder = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.333), Inches(2)
        )
        tf = placeholder.text_frame
        p = tf.paragraphs[0]
        p.text = f"[Image: {image_path}]"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(128, 128, 128)
        p.alignment = PP_ALIGN.CENTER
        
        p = tf.add_paragraph()
        p.text = f"Error: {str(error)}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(192, 192, 192)
        p.alignment = PP_ALIGN.CENTER


def generate_ppt(ppt_structure: dict | str, output_path: str | Path, theme: str = "tech_blue") -> str:
    """
    Convenience function to generate PPT.

    Args:
        ppt_structure: JSON structure dict or path to JSON file
        output_path: Output path for the .pptx file
        theme: Theme name (tech_blue, business_gray, minimal_white, nature_green)

    Returns:
        Path to generated PPT file
    """
    generator = PPTGenerator(theme=theme)
    return generator.generate(ppt_structure, output_path)